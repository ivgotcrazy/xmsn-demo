"""文档解析：PDF / PPT / Word → 按页文本（T2.2 预览与 T2.3 厂商解析复用）。

实现以《厂商解析详细设计》3 章为准：PDF 按页（pypdf）、PPT 按页（python-pptx）、
Word 按段落（python-docx，无页概念整体一页）；单文件解析失败不阻断（返回空）。
"""
from __future__ import annotations

from io import BytesIO


def extract_pages(data: bytes, filename: str) -> list[str]:
    """按扩展名提取逐页文本；无法解析返回空列表（调用方兜底）。"""
    lower = (filename or "").lower()
    try:
        if lower.endswith(".pdf"):
            return _pdf_pages(data)
        if lower.endswith((".ppt", ".pptx")):
            return _ppt_pages(data)
        if lower.endswith((".doc", ".docx")):
            return _doc_pages(data)
    except Exception:  # noqa: BLE001 - 单文件解析失败不阻断
        return []
    return []


def _pdf_pages(data: bytes) -> list[str]:
    from pypdf import PdfReader

    reader = PdfReader(BytesIO(data))
    return [(page.extract_text() or "").strip() for page in reader.pages]


def _ppt_pages(data: bytes) -> list[str]:
    from pptx import Presentation

    prs = Presentation(BytesIO(data))
    pages: list[str] = []
    for slide in prs.slides:
        texts = [
            shape.text.strip()
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text and shape.text.strip()
        ]
        pages.append("\n".join(texts))
    return pages


def _doc_pages(data: bytes) -> list[str]:
    from docx import Document

    doc = Document(BytesIO(data))
    text = "\n".join(p.text.strip() for p in doc.paragraphs if p.text and p.text.strip())
    return [text]
